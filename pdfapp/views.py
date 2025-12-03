"""
PDF app views.
Document upload, conversion, and viewing.
"""

import os
import io
import subprocess
import tempfile
from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth.decorators import login_required
from django.contrib.auth.mixins import LoginRequiredMixin
from django.contrib import messages
from django.views.generic import ListView, CreateView, UpdateView, DeleteView, DetailView
from django.urls import reverse_lazy
from django.http import FileResponse, Http404, HttpResponse
from django.utils import timezone
from django.conf import settings
from django.core.files.base import ContentFile

from .models import PDFDocument, ConversionLog
from .forms import PDFDocumentForm


def convert_docx_to_pdf(docx_path, output_dir):
    """Convert DOCX to PDF using LibreOffice."""
    try:
        result = subprocess.run([
            'soffice',
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', output_dir,
            docx_path
        ], capture_output=True, text=True, timeout=120)
        
        if result.returncode == 0:
            base_name = os.path.splitext(os.path.basename(docx_path))[0]
            pdf_path = os.path.join(output_dir, f"{base_name}.pdf")
            if os.path.exists(pdf_path):
                return pdf_path, None
            return None, "PDF file not generated"
        return None, result.stderr
    except subprocess.TimeoutExpired:
        return None, "Conversion timed out"
    except Exception as e:
        return None, str(e)


def convert_pptx_to_pdf(pptx_path, output_dir):
    """Convert PPTX to PDF using LibreOffice."""
    return convert_docx_to_pdf(pptx_path, output_dir)


def convert_docx_to_html(docx_path):
    """Convert DOCX to HTML preserving styles."""
    try:
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        
        doc = Document(docx_path)
        html_parts = []
        
        # CSS styles
        css = """
        <style>
            body { font-family: 'Calibri', sans-serif; max-width: 800px; margin: 0 auto; padding: 20px; }
            h1 { font-size: 24pt; margin: 20px 0; }
            h2 { font-size: 18pt; margin: 16px 0; }
            h3 { font-size: 14pt; margin: 12px 0; }
            p { margin: 10px 0; line-height: 1.5; }
            .center { text-align: center; }
            .right { text-align: right; }
            .justify { text-align: justify; }
            .bold { font-weight: bold; }
            .italic { font-style: italic; }
            .underline { text-decoration: underline; }
            table { border-collapse: collapse; width: 100%; margin: 10px 0; }
            td, th { border: 1px solid #ddd; padding: 8px; }
            ul, ol { margin: 10px 0; padding-left: 30px; }
        </style>
        """
        
        html_parts.append('<!DOCTYPE html>')
        html_parts.append('<html lang="en">')
        html_parts.append('<head>')
        html_parts.append('<meta charset="UTF-8">')
        html_parts.append('<meta name="viewport" content="width=device-width, initial-scale=1.0">')
        html_parts.append('<title>Converted Document</title>')
        html_parts.append(css)
        html_parts.append('</head>')
        html_parts.append('<body>')
        
        for para in doc.paragraphs:
            if not para.text.strip():
                html_parts.append('<p>&nbsp;</p>')
                continue
            
            # Determine style
            style_name = para.style.name.lower() if para.style else ''
            tag = 'p'
            classes = []
            
            if 'heading 1' in style_name:
                tag = 'h1'
            elif 'heading 2' in style_name:
                tag = 'h2'
            elif 'heading 3' in style_name:
                tag = 'h3'
            
            # Check alignment
            if para.alignment == WD_ALIGN_PARAGRAPH.CENTER:
                classes.append('center')
            elif para.alignment == WD_ALIGN_PARAGRAPH.RIGHT:
                classes.append('right')
            elif para.alignment == WD_ALIGN_PARAGRAPH.JUSTIFY:
                classes.append('justify')
            
            # Build paragraph HTML with run formatting
            content = []
            for run in para.runs:
                text = run.text
                if not text:
                    continue
                
                if run.bold:
                    text = f'<strong>{text}</strong>'
                if run.italic:
                    text = f'<em>{text}</em>'
                if run.underline:
                    text = f'<u>{text}</u>'
                
                content.append(text)
            
            class_attr = f' class="{" ".join(classes)}"' if classes else ''
            html_parts.append(f'<{tag}{class_attr}>{"".join(content)}</{tag}>')
        
        # Process tables
        for table in doc.tables:
            html_parts.append('<table>')
            for row in table.rows:
                html_parts.append('<tr>')
                for cell in row.cells:
                    html_parts.append(f'<td>{cell.text}</td>')
                html_parts.append('</tr>')
            html_parts.append('</table>')
        
        html_parts.append('</body>')
        html_parts.append('</html>')
        
        return '\n'.join(html_parts), None
        
    except Exception as e:
        return None, str(e)


def convert_pptx_to_html(pptx_path):
    """Convert PPTX to HTML preserving styles."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        
        prs = Presentation(pptx_path)
        html_parts = []
        
        # CSS styles
        css = """
        <style>
            body { font-family: 'Calibri', sans-serif; max-width: 960px; margin: 0 auto; padding: 20px; background: #f5f5f5; }
            .slide { background: white; margin: 20px 0; padding: 40px; box-shadow: 0 2px 8px rgba(0,0,0,0.1); min-height: 500px; position: relative; }
            .slide-number { position: absolute; bottom: 10px; right: 20px; color: #999; font-size: 12px; }
            .slide-title { font-size: 28pt; font-weight: bold; margin-bottom: 20px; color: #333; }
            .slide-content { font-size: 14pt; line-height: 1.6; }
            .slide-content ul { margin: 15px 0; padding-left: 30px; }
            .slide-content li { margin: 8px 0; }
            table { border-collapse: collapse; width: 100%; margin: 15px 0; }
            td, th { border: 1px solid #ddd; padding: 10px; text-align: left; }
            th { background: #f0f0f0; }
            .center { text-align: center; }
        </style>
        """
        
        html_parts.append('<!DOCTYPE html>')
        html_parts.append('<html lang="en">')
        html_parts.append('<head>')
        html_parts.append('<meta charset="UTF-8">')
        html_parts.append('<meta name="viewport" content="width=device-width, initial-scale=1.0">')
        html_parts.append('<title>Converted Presentation</title>')
        html_parts.append(css)
        html_parts.append('</head>')
        html_parts.append('<body>')
        
        for slide_num, slide in enumerate(prs.slides, 1):
            html_parts.append(f'<div class="slide">')
            html_parts.append(f'<div class="slide-number">Slide {slide_num}</div>')
            
            for shape in slide.shapes:
                if shape.has_text_frame:
                    is_title = shape == slide.shapes.title if hasattr(slide.shapes, 'title') else False
                    
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if not text:
                            continue
                        
                        if is_title:
                            html_parts.append(f'<div class="slide-title">{text}</div>')
                            is_title = False
                        else:
                            # Check if it's a bullet point
                            if paragraph.level > 0:
                                html_parts.append(f'<ul><li>{text}</li></ul>')
                            else:
                                html_parts.append(f'<div class="slide-content">{text}</div>')
                
                elif shape.has_table:
                    table = shape.table
                    html_parts.append('<table>')
                    for row in table.rows:
                        html_parts.append('<tr>')
                        for cell in row.cells:
                            html_parts.append(f'<td>{cell.text}</td>')
                        html_parts.append('</tr>')
                    html_parts.append('</table>')
            
            html_parts.append('</div>')
        
        html_parts.append('</body>')
        html_parts.append('</html>')
        
        return '\n'.join(html_parts), None
        
    except Exception as e:
        return None, str(e)


def process_document(document):
    """Process uploaded document - convert to PDF and HTML."""
    
    document.status = 'processing'
    document.save()
    
    ConversionLog.objects.create(
        document=document,
        action='start',
        message='Starting document conversion'
    )
    
    try:
        source_path = document.source_file.path
        
        # Create temp directory for output
        with tempfile.TemporaryDirectory() as temp_dir:
            # Convert to PDF
            if document.is_docx:
                pdf_path, error = convert_docx_to_pdf(source_path, temp_dir)
            else:
                pdf_path, error = convert_pptx_to_pdf(source_path, temp_dir)
            
            if pdf_path and os.path.exists(pdf_path):
                with open(pdf_path, 'rb') as f:
                    pdf_name = f"{os.path.splitext(os.path.basename(source_path))[0]}.pdf"
                    document.pdf_file.save(pdf_name, ContentFile(f.read()), save=False)
                
                ConversionLog.objects.create(
                    document=document,
                    action='pdf_success',
                    message='PDF conversion successful'
                )
            elif error:
                ConversionLog.objects.create(
                    document=document,
                    action='pdf_error',
                    message=f'PDF conversion failed: {error}'
                )
            
            # Convert to HTML
            if document.is_docx:
                html_content, error = convert_docx_to_html(source_path)
            else:
                html_content, error = convert_pptx_to_html(source_path)
            
            if html_content:
                document.html_content = html_content
                
                # Also save as file
                html_name = f"{os.path.splitext(os.path.basename(source_path))[0]}.html"
                document.html_file.save(html_name, ContentFile(html_content.encode()), save=False)
                
                ConversionLog.objects.create(
                    document=document,
                    action='html_success',
                    message='HTML conversion successful'
                )
            elif error:
                ConversionLog.objects.create(
                    document=document,
                    action='html_error',
                    message=f'HTML conversion failed: {error}'
                )
        
        document.status = 'completed'
        document.processed_at = timezone.now()
        document.save()
        
        ConversionLog.objects.create(
            document=document,
            action='complete',
            message='Document processing completed'
        )
        
        return True
        
    except Exception as e:
        document.status = 'failed'
        document.error_message = str(e)
        document.save()
        
        ConversionLog.objects.create(
            document=document,
            action='error',
            message=f'Processing failed: {str(e)}'
        )
        
        return False


# Dashboard Views
class PDFDocumentListView(LoginRequiredMixin, ListView):
    """List user's PDF documents."""
    
    model = PDFDocument
    template_name = 'pdfapp/document_list.html'
    context_object_name = 'documents'
    
    def get_queryset(self):
        return PDFDocument.objects.filter(user=self.request.user)
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        context['module'] = 'pdf'
        return context


class PDFDocumentCreateView(LoginRequiredMixin, CreateView):
    """Upload and convert document."""
    
    model = PDFDocument
    form_class = PDFDocumentForm
    template_name = 'pdfapp/document_form.html'
    success_url = reverse_lazy('pdfapp:list')
    
    def form_valid(self, form):
        form.instance.user = self.request.user
        response = super().form_valid(form)
        
        # Process the document
        process_document(self.object)
        
        if self.object.status == 'completed':
            messages.success(self.request, 'Document converted successfully!')
        else:
            messages.warning(self.request, f'Conversion completed with issues: {self.object.error_message}')
        
        return response
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        context['module'] = 'pdf'
        context['title'] = 'Upload Document'
        return context


class PDFDocumentDetailView(LoginRequiredMixin, DetailView):
    """View document details and preview."""
    
    model = PDFDocument
    template_name = 'pdfapp/document_detail.html'
    context_object_name = 'document'
    
    def get_queryset(self):
        return PDFDocument.objects.filter(user=self.request.user)
    
    def get_object(self, queryset=None):
        obj = super().get_object(queryset)
        obj.view_count += 1
        obj.save()
        return obj
    
    def get_context_data(self, **kwargs):
        context = super().get_context_data(**kwargs)
        context['module'] = 'pdf'
        context['logs'] = self.object.logs.all()[:10]
        return context


class PDFDocumentDeleteView(LoginRequiredMixin, DeleteView):
    """Delete document."""
    
    model = PDFDocument
    template_name = 'dashboard/confirm_delete.html'
    success_url = reverse_lazy('pdfapp:list')
    
    def get_queryset(self):
        return PDFDocument.objects.filter(user=self.request.user)
    
    def form_valid(self, form):
        messages.success(self.request, 'Document deleted successfully!')
        return super().form_valid(form)


@login_required
def view_pdf(request, pk):
    """View PDF in browser."""
    
    document = get_object_or_404(PDFDocument, pk=pk, user=request.user)
    
    if not document.pdf_file:
        raise Http404("PDF not available")
    
    try:
        response = FileResponse(
            document.pdf_file.open('rb'),
            content_type='application/pdf'
        )
        response['Content-Disposition'] = f'inline; filename="{os.path.basename(document.pdf_file.name)}"'
        return response
    except Exception as e:
        raise Http404("File not found")


@login_required
def view_html(request, pk):
    """View HTML version."""
    
    document = get_object_or_404(PDFDocument, pk=pk, user=request.user)
    
    if document.html_content:
        return HttpResponse(document.html_content, content_type='text/html')
    elif document.html_file:
        try:
            return FileResponse(
                document.html_file.open('rb'),
                content_type='text/html'
            )
        except:
            pass
    
    raise Http404("HTML not available")


@login_required
def export_html(request, pk):
    """Export HTML as downloadable file."""
    
    document = get_object_or_404(PDFDocument, pk=pk, user=request.user)
    
    if document.html_content:
        response = HttpResponse(document.html_content, content_type='text/html')
        filename = f"{os.path.splitext(os.path.basename(document.source_file.name))[0]}.html"
        response['Content-Disposition'] = f'attachment; filename="{filename}"'
        return response
    elif document.html_file:
        try:
            response = FileResponse(
                document.html_file.open('rb'),
                as_attachment=True,
                filename=os.path.basename(document.html_file.name)
            )
            return response
        except:
            pass
    
    raise Http404("HTML not available")


@login_required
def reprocess_document(request, pk):
    """Reprocess a document."""
    
    document = get_object_or_404(PDFDocument, pk=pk, user=request.user)
    
    process_document(document)
    
    if document.status == 'completed':
        messages.success(request, 'Document reprocessed successfully!')
    else:
        messages.error(request, f'Reprocessing failed: {document.error_message}')
    
    return redirect('pdfapp:detail', pk=pk)
